"""Generate concise 12-slide TalentIQ overview with visual layouts."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parent.parent / "docs" / "TalentIQ-Project-Overview-Advanced.pptx"

# Brand
NAVY = RGBColor(0x0F, 0x17, 0x2A)
SLATE = RGBColor(0x1E, 0x29, 0x3B)
BLUE = RGBColor(0x25, 0x63, 0xEB)
BLUE_DK = RGBColor(0x1D, 0x4E, 0xD8)
BLUE_LT = RGBColor(0xDB, 0xEA, 0xFE)
TEAL = RGBColor(0x0D, 0x94, 0x88)
TEAL_LT = RGBColor(0xCC, 0xFB, 0xF1)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_LT = RGBColor(0xED, 0xE9, 0xFE)
AMBER = RGBColor(0xD9, 0x77, 0x06)
AMBER_LT = RGBColor(0xFF, 0xFB, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x64, 0x74, 0x8B)
TEXT = RGBColor(0x1E, 0x29, 0x3B)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
HDR = RGBColor(0x1E, 0x40, 0xAF)
ROW_ALT = RGBColor(0xF8, 0xFA, 0xFC)

SW = Inches(13.333)
SH = Inches(7.5)

CARD_PALETTE = [
    (BLUE, BLUE_LT),
    (TEAL, TEAL_LT),
    (PURPLE, PURPLE_LT),
    (AMBER, AMBER_LT),
    (BLUE_DK, BLUE_LT),
    (SLATE, RGBColor(0xF1, 0xF5, 0xF9)),
]


def slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(s, color=WHITE):
    f = s.background.fill
    f.solid()
    f.fore_color.rgb = color


def bar(s, y=Inches(0), h=Inches(0.06), color=BLUE):
    b = s.shapes.add_shape(1, Inches(0), y, SW, h)
    b.fill.solid()
    b.fill.fore_color.rgb = color
    b.line.fill.background()


def heading(s, text, y=Inches(0.38), size=30, color=NAVY):
    box = s.shapes.add_textbox(Inches(0.6), y, Inches(12.2), Inches(0.65))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color


def sub(s, text, y=Inches(0.98)):
    box = s.shapes.add_textbox(Inches(0.6), y, Inches(12.2), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(13)
    p.font.color.rgb = MUTED


def card(s, x, y, w, h, title, body, accent=BLUE, fill=BLUE_LT, title_size=13, body_size=10):
    shape = s.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(title_size)
    p0.font.bold = True
    p0.font.color.rgb = accent
    p1 = tf.add_paragraph()
    p1.text = body
    p1.font.size = Pt(body_size)
    p1.font.color.rgb = TEXT
    p1.space_before = Pt(4)


def stat_card(s, x, y, w, h, value, label, accent=BLUE, fill=BLUE_LT):
    shape = s.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = BORDER
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = value
    p0.font.size = Pt(32)
    p0.font.bold = True
    p0.font.color.rgb = accent
    p0.alignment = PP_ALIGN.CENTER
    p1 = tf.add_paragraph()
    p1.text = label
    p1.font.size = Pt(11)
    p1.font.bold = True
    p1.font.color.rgb = NAVY
    p1.alignment = PP_ALIGN.CENTER


def table(s, headers, rows, top, col_widths=None, row_h=0.38):
    nr, nc = len(rows) + 1, len(headers)
    ts = s.shapes.add_table(nr, nc, Inches(0.55), top, Inches(12.2), Inches(row_h * nr))
    t = ts.table
    if col_widths:
        for i, w in enumerate(col_widths):
            t.columns[i].width = Inches(w)
    for ci, h in enumerate(headers):
        c = t.cell(0, ci)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = HDR
        for p in c.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(10)
            p.font.color.rgb = WHITE
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = t.cell(ri + 1, ci)
            c.text = val
            if ri % 2:
                c.fill.solid()
                c.fill.fore_color.rgb = ROW_ALT
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = TEXT


# ─── SLIDES ──────────────────────────────────────────────────────

def s01_title(prs):
    s = slide(prs)
    bg(s, NAVY)
    bar(s, Inches(5.1), Inches(0.07), BLUE)
    t = s.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(12), Inches(1.1))
    t.text_frame.paragraphs[0].text = "TalentIQ"
    t.text_frame.paragraphs[0].font.size = Pt(54)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE
    sub2 = s.shapes.add_textbox(Inches(0.7), Inches(3.15), Inches(12), Inches(0.6))
    sub2.text_frame.paragraphs[0].text = "AI-Powered Workforce Intelligence Platform"
    sub2.text_frame.paragraphs[0].font.size = Pt(24)
    sub2.text_frame.paragraphs[0].font.color.rgb = BLUE_LT
    m = s.shapes.add_textbox(Inches(0.7), Inches(4.2), Inches(12), Inches(0.5))
    m.text_frame.paragraphs[0].text = "TalentIQ-RGM-Webpage  •  Enterprise HR, Learning & Talent Management"
    m.text_frame.paragraphs[0].font.size = Pt(13)
    m.text_frame.paragraphs[0].font.color.rgb = MUTED


def s02_overview(prs):
    """Stats cards + mission statement."""
    s = slide(prs)
    bg(s)
    bar(s)
    heading(s, "Platform Overview")
    sub(s, "Single enterprise system connecting skills, learning, analytics & AI")

    stats = [("40+", "Screens"), ("168", "APIs"), ("50+", "DB Models"), ("8", "Roles"), ("45+", "Services"), ("9", "Core Modules")]
    cw, gap = Inches(1.85), Inches(0.18)
    x0 = Inches(0.55)
    for i, (v, l) in enumerate(stats):
        accent, fill = CARD_PALETTE[i % len(CARD_PALETTE)]
        stat_card(s, x0 + i * (cw + gap), Inches(1.55), cw, Inches(1.55), v, l, accent, fill)

    card(s, Inches(0.55), Inches(3.35), Inches(12.2), Inches(1.45),
         "Mission",
         "Unify workforce data — skills, training, certifications, and performance — into one live platform "
         "so HR, managers, and executives can make faster, data-driven talent decisions.",
         BLUE, BLUE_LT, 14, 12)

    cards = [
        ("Skill Matrix", "Heatmap readiness by employee & department"),
        ("Learning (LMS)", "Courses, assessments, certifications"),
        ("Career & Promotion", "Role paths & readiness scoring"),
        ("AI Copilot", "Natural-language workforce queries"),
    ]
    cw2 = Inches(2.95)
    for i, (t, b) in enumerate(cards):
        accent, fill = CARD_PALETTE[i]
        card(s, Inches(0.55) + i * (cw2 + Inches(0.15)), Inches(5.05), cw2, Inches(1.15), t, b, accent, fill, 12, 10)


def s03_journey(prs):
    """Visual lifecycle flow."""
    s = slide(prs)
    bg(s)
    bar(s)
    heading(s, "Core Employee Lifecycle")
    sub(s, "Every stage connected through live PostgreSQL data")

    steps = [
        ("1", "Employee", "Onboard & profile"),
        ("2", "Role & Level", "Job role + experience"),
        ("3", "Skills", "Required competencies"),
        ("4", "Learning", "Courses & roadmaps"),
        ("5", "Assessment", "Validate knowledge"),
        ("6", "Certification", "Issue credentials"),
        ("7", "Promotion", "Readiness score"),
    ]
    bw, bh, gap = Inches(1.55), Inches(1.65), Inches(0.14)
    x0 = Inches(0.4)
    y = Inches(1.65)
    accents = [BLUE, TEAL, PURPLE, BLUE_DK, TEAL, AMBER, BLUE]

    for i, (num, title, desc) in enumerate(steps):
        x = x0 + i * (bw + gap)
        shape = s.shapes.add_shape(1, x, y, bw, bh)
        accent = accents[i]
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = accent
        shape.line.width = Pt(2)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.08)
        p0 = tf.paragraphs[0]
        p0.text = num
        p0.font.size = Pt(20)
        p0.font.bold = True
        p0.font.color.rgb = accent
        p0.alignment = PP_ALIGN.CENTER
        p1 = tf.add_paragraph()
        p1.text = title
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = NAVY
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(9)
        p2.font.color.rgb = MUTED
        p2.alignment = PP_ALIGN.CENTER
        if i < len(steps) - 1:
            arr = s.shapes.add_textbox(x + bw, y + Inches(0.6), gap, Inches(0.3))
            arr.text_frame.paragraphs[0].text = "▶"
            arr.text_frame.paragraphs[0].font.size = Pt(12)
            arr.text_frame.paragraphs[0].font.color.rgb = MUTED
            arr.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    layers = [
        ("Organization", "Departments • Teams • Org Chart • Employees", TEAL, TEAL_LT),
        ("Skills", "Framework • Matrix • Competency • Validation", PURPLE, PURPLE_LT),
        ("Learning", "Courses • Assessments • Certs • xAPI/LRS", BLUE, BLUE_LT),
        ("Intelligence", "Analytics • AI Copilot • War Room • Forecasting", NAVY, RGBColor(0xE2, 0xE8, 0xF0)),
    ]
    for i, (title, desc, accent, fill) in enumerate(layers):
        card(s, Inches(0.55) + i * Inches(3.15), Inches(3.65), Inches(3.0), Inches(1.2), title, desc, accent, fill, 13, 10)


def s04_modules(prs):
    """7-section module table."""
    s = slide(prs)
    bg(s)
    bar(s)
    heading(s, "Platform Modules — 7 Functional Areas", size=28)
    sub(s, "40+ screens with role-based sidebar navigation")
    table(s,
          ["Section", "Key Screens", "What It Does"],
          [
              ["Command Center", "Dashboard, AI Copilot, Executive War Room", "KPIs, AI queries, C-suite readiness view"],
              ["Workforce", "Employees, Teams, Departments, Org Chart, Talent Marketplace", "People management & internal mobility"],
              ["Skill Intelligence", "Skill Matrix, Framework, Competency Models, Validation", "Taxonomy, gap heatmaps, role requirements"],
              ["Learning & Dev", "Programs, Pathways, Courses, Assessments, Certifications", "Full LMS with xAPI progress tracking"],
              ["Performance", "Career, Goals, Reviews, Promotion, Succession", "Career paths, reviews, bench strength"],
              ["Analytics", "Workforce, Executive, Dept, Team, Compliance, Forecasting, Reports", "Multi-level dashboards + PDF/Excel export"],
              ["Administration", "Users, RBAC, Audit Logs, Screen Access, Settings", "Governance, security & system config"],
          ],
          Inches(1.35), [2.0, 4.5, 5.7], 0.42)


def s05_command_ai(prs):
    """AI Copilot cards + command center."""
    s = slide(prs)
    bg(s)
    bar(s)
    heading(s, "Command Center & AI Workforce Copilot", size=26)
    sub(s, "Intent-driven intelligence over live database — not a generic chatbot")

    card(s, Inches(0.55), Inches(1.4), Inches(3.8), Inches(2.0),
         "Dashboard", "Workforce KPIs, alerts, quick actions, notification stream", BLUE, BLUE_LT, 13, 11)
    card(s, Inches(4.55), Inches(1.4), Inches(3.8), Inches(2.0),
         "Executive War Room", "Org health, readiness risk, critical role coverage for CEO/Admin", PURPLE, PURPLE_LT, 13, 11)
    card(s, Inches(8.55), Inches(1.4), Inches(4.2), Inches(2.0),
         "AI Copilot Pipeline", "Query → Intent classify → RBAC scope → Live DB → Ranked results + drill-down", TEAL, TEAL_LT, 13, 11)

    intents = [
        ("Promotion Ready", "Eligible employees"),
        ("Skill Gaps", "Matrix & missing skills"),
        ("Cert Risk", "Expiring credentials"),
        ("Compliance", "Non-compliant segments"),
        ("Learning", "Training progress"),
        ("Succession", "Bench strength"),
        ("Attrition", "At-risk employees"),
        ("Dept Analysis", "Cross-team compare"),
        ("Performers", "Top/bottom ranking"),
        ("Org Health", "Overall readiness"),
    ]
    iw, ih = Inches(1.15), Inches(0.72)
    for i, (t, b) in enumerate(intents):
        col, row = i % 5, i // 5
        x = Inches(0.55) + col * (iw + Inches(0.12))
        y = Inches(3.65) + row * (ih + Inches(0.1))
        accent, fill = CARD_PALETTE[i % len(CARD_PALETTE)]
        card(s, x, y, iw, ih, t, b, accent, fill, 9, 8)


def s06_workforce_skills(prs):
    """Two-column feature cards."""
    s = slide(prs)
    bg(s)
    bar(s)
    heading(s, "Workforce & Skill Intelligence", size=26)
    sub(s, "Organizational structure meets competency management")

    wf = [
        ("Employees & Contractors", "Profiles, manager hierarchy, skill & learning history"),
        ("Departments & Teams", "Hierarchical org structure with team leaders"),
        ("Organization Chart", "Interactive visual hierarchy"),
        ("Talent Marketplace", "Internal mobility opportunities & applications"),
    ]
    sk = [
        ("Skill Matrix", "Heatmap: employee × skill × department readiness"),
        ("Skill Framework", "Categories, levels, taxonomy & relations"),
        ("Competency Models", "Role + experience level → required skills"),
        ("Skill Validation", "Manager verification of proficiencies"),
    ]
    for i, (t, b) in enumerate(wf):
        card(s, Inches(0.55), Inches(1.45) + i * Inches(1.22), Inches(5.9), Inches(1.05), t, b, TEAL, TEAL_LT, 12, 10)
    for i, (t, b) in enumerate(sk):
        card(s, Inches(6.75), Inches(1.45) + i * Inches(1.22), Inches(5.9), Inches(1.05), t, b, PURPLE, PURPLE_LT, 12, 10)

    # connector arrow label
    mid = s.shapes.add_textbox(Inches(6.35), Inches(3.2), Inches(0.5), Inches(0.4))
    mid.text_frame.paragraphs[0].text = "↔"
    mid.text_frame.paragraphs[0].font.size = Pt(22)
    mid.text_frame.paragraphs[0].font.color.rgb = MUTED


def s07_learning_performance(prs):
    """L&D + Performance cards in grid."""
    s = slide(prs)
    bg(s)
    bar(s)
    heading(s, "Learning, Performance & Career", size=26)
    sub(s, "Full LMS + career progression pipeline")

    items = [
        ("Courses & Programs", "Module → Lesson hierarchy\nVideo, PDF, Quiz, Assignment", BLUE, BLUE_LT),
        ("Assessments", "Question bank, timed attempts\nPass rates & grading", TEAL, TEAL_LT),
        ("Certifications", "Templates, expiry tracking\nPublic verify URL", PURPLE, PURPLE_LT),
        ("xAPI / LRS", "Learning Record Store\nExternal platform sync", AMBER, AMBER_LT),
        ("Career Progression", "Role paths & experience levels\nPromotion targets", BLUE_DK, BLUE_LT),
        ("Goals & Reviews", "Performance cycles\nTeam & dept reviews", TEAL, TEAL_LT),
        ("Promotion Readiness", "Skill coverage vs role requirements\nAutomated scoring", PURPLE, PURPLE_LT),
        ("Succession Planning", "Critical roles, successors\nBench strength analysis", NAVY, RGBColor(0xE2, 0xE8, 0xF0)),
    ]
    cw, ch = Inches(2.95), Inches(1.55)
    for i, (t, b, ac, fi) in enumerate(items):
        col, row = i % 4, i // 4
        card(s, Inches(0.55) + col * (cw + Inches(0.12)), Inches(1.45) + row * (ch + Inches(0.12)), cw, ch, t, b, ac, fi, 11, 9)


def s08_analytics_admin(prs):
    """Analytics table + admin cards."""
    s = slide(prs)
    bg(s)
    bar(s)
    heading(s, "Analytics, Reports & Administration", size=26)
    sub(s, "Scoped dashboards by role + enterprise governance")

    table(s,
          ["Analytics View", "Audience", "Key Metrics"],
          [
              ["Workforce Analytics", "HR / All managers", "Readiness, compliance, benchmarks"],
              ["Executive Analytics", "CEO, Admin, HR", "Org-wide KPIs & trends"],
              ["Department / Team", "Dept Manager, Team Lead", "Scoped skill & learning metrics"],
              ["Compliance", "HR Manager", "Mandatory training & cert status"],
              ["Forecasting", "Executive", "Predictive talent demand"],
              ["Reports", "All authorized", "PDF & Excel export"],
          ],
          Inches(1.35), [2.8, 2.5, 6.9], 0.36)

    admin = [
        ("Users & RBAC", "Role-permission matrix, screen overrides"),
        ("Audit Logs", "Immutable trail — login, CRUD, exports"),
        ("Access Requests", "Approve/deny permission elevation"),
        ("System Settings", "Security, email, notifications, appearance"),
    ]
    for i, (t, b) in enumerate(admin):
        card(s, Inches(0.55) + i * Inches(3.08), Inches(4.15), Inches(2.95), Inches(0.95), t, b, SLATE, RGBColor(0xF1, 0xF5, 0xF9), 11, 9)


def s09_data_api(prs):
    s = slide(prs)
    bg(s)
    bar(s)
    heading(s, "Data Model & API Layer", size=26)
    sub(s, "PostgreSQL + Prisma ORM + 168 REST endpoints")

    for i, (v, l, ac, fi) in enumerate([
        ("50+", "DB Models", BLUE, BLUE_LT),
        ("168", "API Routes", TEAL, TEAL_LT),
        ("40+", "Permissions", PURPLE, PURPLE_LT),
        ("Soft", "Delete", AMBER, AMBER_LT),
    ]):
        stat_card(s, Inches(0.55) + i * Inches(2.15), Inches(1.35), Inches(1.95), Inches(1.15), v, l, ac, fi)

    domains = [
        ("RBAC & Access", "Role, Permission, Screen, ScreenAccess, AccessRequest"),
        ("Organization", "Department, Team, User, JobRole, ExperienceLevel"),
        ("Skills", "Skill, EmployeeSkill, RoleSkillRequirement, Mappings"),
        ("Learning & LRS", "Course, Enrollment, Roadmap, OpenCourse, Progress"),
        ("Assessments & Certs", "QuestionBank, Assessment, Attempt, Certificate"),
        ("Intelligence", "Forecast, Succession, TalentOpportunity, AuditLog"),
    ]
    for i, (t, b) in enumerate(domains):
        col, row = i % 3, i // 3
        accent, fill = CARD_PALETTE[i]
        card(s, Inches(0.55) + col * Inches(4.05), Inches(2.75) + row * Inches(1.35), Inches(3.85), Inches(1.2), t, b, accent, fill, 11, 9)

    api = s.shapes.add_shape(1, Inches(0.55), Inches(5.55), Inches(12.2), Inches(0.85))
    api.fill.solid()
    api.fill.fore_color.rgb = NAVY
    api.line.fill.background()
    tf = api.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = "API: Auth • Users • RBAC • Skills • Courses • Learning/LRS • Assessments • Certs • Analytics • AI Copilot • War Room • Forecasting • Reports • Health"
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE


def s10_rbac(prs):
    """Role table + security cards."""
    s = slide(prs)
    bg(s)
    bar(s)
    heading(s, "RBAC & Security", size=26)
    sub(s, "8 roles • 40+ permissions • screen-level access control")

    table(s,
          ["Role", "Data Scope", "Primary Access"],
          [
              ["Employee", "Self", "Learning, courses, assessments, career"],
              ["Instructor", "Self + teach", "+ Learning programs"],
              ["Team Leader", "Team", "+ Employees, skill matrix, team analytics"],
              ["Dept Manager", "Department", "+ Teams, dept analytics, succession"],
              ["HR Manager", "Organization", "Full HR, skills, L&D, compliance"],
              ["CEO", "Executive", "War room, exec analytics, marketplace"],
              ["Admin / Owner", "System", "All screens, RBAC, audit, settings"],
          ],
          Inches(1.35), [1.8, 1.8, 8.6], 0.34)

    sec = [
        ("JWT Auth", "15m access + 7d refresh tokens"),
        ("CSRF", "Double-submit cookie protection"),
        ("Rate Limit", "Redis-backed API throttling"),
        ("Audit Log", "Append-only activity trail"),
    ]
    for i, (t, b) in enumerate(sec):
        card(s, Inches(0.55) + i * Inches(3.08), Inches(4.35), Inches(2.95), Inches(0.9), t, b, BLUE_DK, BLUE_LT, 11, 9)


def s11_tech(prs):
    """Tech stack two-column cards."""
    s = slide(prs)
    bg(s)
    bar(s)
    heading(s, "Technology Stack", size=26)
    sub(s, "Modern full-stack monorepo — Next.js serves frontend & backend")

    fe = [
        ("Frontend", "Next.js 15 • React 19 • TypeScript • Tailwind 4 • Shadcn UI"),
        ("State & Data", "TanStack Query • Zustand • React Hook Form • Zod"),
        ("UI & Charts", "Framer Motion • Recharts • TanStack Table • Lucide"),
    ]
    be = [
        ("Backend", "Next.js Route Handlers • 45+ domain services • Prisma ORM"),
        ("Infrastructure", "PostgreSQL 16 • Redis 7 • AWS S3 • Sentry • Pino logs"),
        ("Integrations", "xAPI/LRS • Zoho People • Email • PDF/Excel exports"),
    ]
    for i, (t, b) in enumerate(fe):
        card(s, Inches(0.55), Inches(1.45) + i * Inches(1.35), Inches(5.9), Inches(1.15), t, b, BLUE, BLUE_LT, 13, 11)
    for i, (t, b) in enumerate(be):
        card(s, Inches(6.75), Inches(1.45) + i * Inches(1.35), Inches(5.9), Inches(1.15), t, b, TEAL, TEAL_LT, 13, 11)

    folders = s.shapes.add_shape(1, Inches(0.55), Inches(5.55), Inches(12.2), Inches(1.0))
    folders.fill.solid()
    folders.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)
    folders.line.color.rgb = BORDER
    tf = folders.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = "Structure:  app/ (pages + 168 APIs)  •  components/  •  services/  •  lib/  •  prisma/  •  config/  •  docs/  •  __tests__/"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT


def s12_devops(prs):
    """Deployment cards + QA."""
    s = slide(prs)
    bg(s)
    bar(s)
    heading(s, "DevOps, Quality & Deployment", size=26)
    sub(s, "Production-ready with CI/CD and operational checklists")

    deploy = [
        ("Docker Compose", "app + PostgreSQL + Redis\nRecommended for VPS"),
        ("Vercel Cloud", "Managed DB (Neon/Supabase)\n+ Upstash Redis"),
        ("Standalone Node", "npm ci → migrate → build → start"),
    ]
    for i, (t, b) in enumerate(deploy):
        card(s, Inches(0.55) + i * Inches(4.05), Inches(1.45), Inches(3.85), Inches(1.35), t, b, CARD_PALETTE[i][0], CARD_PALETTE[i][1], 12, 10)

    qa = [
        ("Vitest", "RBAC, screen access, user tests"),
        ("Validation Scripts", "rbac:validate, screens:verify"),
        ("GitHub Actions CI", "Lint, build, migrate, audit"),
        ("Health Checks", "/api/health/live & /ready"),
    ]
    for i, (t, b) in enumerate(qa):
        card(s, Inches(0.55) + i * Inches(3.08), Inches(3.05), Inches(2.95), Inches(1.0), t, b, PURPLE, PURPLE_LT, 11, 9)

    card(s, Inches(0.55), Inches(4.3), Inches(12.2), Inches(1.05),
         "Getting Started", "npm run dev (port 3000)  •  npm run db:seed  •  Admin: admin@talentiq.com  •  See .env.example",
         NAVY, RGBColor(0xE2, 0xE8, 0xF0), 12, 11)


def s13_close(prs):
    s = slide(prs)
    bg(s, NAVY)
    bar(s, Inches(5.1), Inches(0.07), BLUE)

    diffs = [
        ("Unified Platform", "Skills + LMS + Analytics + AI in one system"),
        ("Live-Data AI", "Deterministic copilot, not generic chat"),
        ("Granular RBAC", "Screen-level access with audit trail"),
        ("Full Lifecycle", "Onboarding → promotion readiness"),
    ]
    for i, (t, b) in enumerate(diffs):
        card(s, Inches(0.55) + i * Inches(3.08), Inches(1.5), Inches(2.95), Inches(1.15), t, b, BLUE, BLUE_LT, 12, 10)

    t = s.shapes.add_textbox(Inches(0.7), Inches(3.2), Inches(12), Inches(0.9))
    t.text_frame.paragraphs[0].text = "Thank You"
    t.text_frame.paragraphs[0].font.size = Pt(44)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE

    m = s.shapes.add_textbox(Inches(0.7), Inches(4.3), Inches(12), Inches(0.8))
    m.text_frame.paragraphs[0].text = "TalentIQ — Turning Workforce Data into Actionable Talent Decisions"
    m.text_frame.paragraphs[0].font.size = Pt(16)
    m.text_frame.paragraphs[0].font.color.rgb = BLUE_LT


def main():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    s01_title(prs)           # 1
    s02_overview(prs)        # 2  stats + mission cards
    s03_journey(prs)         # 3  flow + layer cards
    s04_modules(prs)         # 4  module table
    s05_command_ai(prs)      # 5  AI cards grid
    s06_workforce_skills(prs)# 6  two-column cards
    s07_learning_performance(prs)  # 7  8-card grid
    s08_analytics_admin(prs) # 8  table + admin cards
    s09_data_api(prs)        # 9  data model cards
    s10_rbac(prs)            # 10 role table + security
    s11_tech(prs)            # 11 tech stack
    s12_devops(prs)          # 12 devops
    s13_close(prs)           # 13 thank you

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(OUTPUT))
    except PermissionError:
        alt = OUTPUT.with_name("TalentIQ-Project-Overview-Final.pptx")
        prs.save(str(alt))
        print(f"Original locked — saved to: {alt} ({len(prs.slides)} slides)")
        return
    print(f"Created: {OUTPUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
